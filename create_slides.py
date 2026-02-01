#!/usr/bin/env python3
"""
Generate NL2SQL MCP Architecture PowerPoint slides - Compact Version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullet_points, font_size=18):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if point.startswith("    "):
            p.level = 2
            p.text = point.strip()
        elif point.startswith("  "):
            p.level = 1
            p.text = point.strip()
        else:
            p.level = 0
            p.text = point
        for run in p.runs:
            run.font.size = Pt(font_size)

    return slide

def add_box(slide, x, y, w, h, text, subtext, color, text_size=10, subtext_size=8):
    """Add a colored box with title and subtitle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(80, 80, 80)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(text_size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtext:
        p2 = tf.add_paragraph()
        p2.text = subtext
        p2.font.size = Pt(subtext_size)
        p2.font.color.rgb = RGBColor(230, 230, 230)
        p2.alignment = PP_ALIGN.CENTER

    return shape

def add_diamond(slide, x, y, w, h, text, color):
    """Add a diamond decision shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(80, 80, 80)

    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x, y, w, h, direction="right", color=RGBColor(100, 100, 100)):
    """Add an arrow"""
    shape_types = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW
    }
    arrow = slide.shapes.add_shape(shape_types.get(direction, MSO_SHAPE.RIGHT_ARROW), x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def add_text(slide, x, y, w, h, text, size=10, bold=False, italic=False, color=RGBColor(0,0,0), align="left"):
    """Add a text box"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.italic = italic
    tf.paragraphs[0].font.color.rgb = color
    if align == "center":
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return tb

def add_flowchart_slide(prs):
    """Compact flowchart that fits on slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, Inches(0.3), Inches(0.15), Inches(6), Inches(0.35),
             "NL2SQL Process Flow", size=22, bold=True)

    # Colors
    BLUE = RGBColor(66, 133, 244)
    GREEN = RGBColor(52, 168, 83)
    YELLOW = RGBColor(200, 160, 0)
    RED = RGBColor(219, 68, 55)
    PURPLE = RGBColor(142, 86, 178)
    TEAL = RGBColor(0, 137, 123)
    GRAY = RGBColor(117, 117, 117)

    # Compact dimensions
    box_w = Inches(1.55)
    box_h = Inches(0.52)
    arrow_w = Inches(0.22)
    arrow_h = Inches(0.15)
    gap = Inches(0.08)
    diamond_w = Inches(0.55)
    diamond_h = Inches(0.55)

    # === ROW 1: Input flow (horizontal) ===
    row1_y = Inches(0.55)
    x = Inches(0.25)

    # 1. User
    add_box(slide, x, row1_y, box_w, box_h, "1. User", "NL Question", BLUE, 9, 7)
    x += box_w + gap
    add_arrow(slide, x, row1_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "right")
    x += arrow_w + gap

    # 2. LLM
    add_box(slide, x, row1_y, box_w, box_h, "2. LLM", "nl_query tool", GREEN, 9, 7)
    x += box_w + gap
    add_arrow(slide, x, row1_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "right")
    x += arrow_w + gap

    # 3. TS MCP
    add_box(slide, x, row1_y, box_w, box_h, "3. TS MCP", "Orchestrate", YELLOW, 9, 7)
    x += box_w + gap
    add_arrow(slide, x, row1_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "right")
    x += arrow_w + gap

    # 4. Python
    python_x = x
    add_box(slide, x, row1_y, box_w, box_h, "4. Python", "Build prompt", RED, 9, 7)
    x += box_w + gap
    add_arrow(slide, x, row1_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "right")
    x += arrow_w + gap

    # 5. Ollama
    add_box(slide, x, row1_y, box_w, box_h, "5. Ollama", "Generate SQL", PURPLE, 9, 7)

    # Down arrow from Ollama
    ollama_center_x = x + box_w/2
    add_arrow(slide, ollama_center_x - arrow_h/2, row1_y + box_h + gap, arrow_h, arrow_w, "down")

    # === ROW 2: Validation flow ===
    row2_y = row1_y + box_h + gap + arrow_w + gap + Inches(0.05)

    # 6. Semantic Validation (under Ollama, flows left)
    sem_x = ollama_center_x - box_w/2
    add_box(slide, sem_x, row2_y, box_w, box_h, "6. Semantic", "Entity check", PURPLE, 9, 7)

    # Arrow left
    add_arrow(slide, sem_x - arrow_w - gap, row2_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "left")

    # 7. Structural Validation
    struct_x = sem_x - arrow_w - gap - gap - box_w
    add_box(slide, struct_x, row2_y, box_w, box_h, "7. Structural", "SELECT-only", YELLOW, 9, 7)

    # Arrow left to decision
    add_arrow(slide, struct_x - arrow_w - gap, row2_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "left")

    # Decision diamond: Valid?
    dec1_x = struct_x - arrow_w - gap - gap - diamond_w
    dec1_y = row2_y + box_h/2 - diamond_h/2
    add_diamond(slide, dec1_x, dec1_y, diamond_w, diamond_h, "OK?", GRAY)

    # Yes arrow down
    add_text(slide, dec1_x - Inches(0.25), dec1_y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.15),
             "Yes", size=7, bold=True, color=GREEN)
    add_arrow(slide, dec1_x + diamond_w/2 - arrow_h/2, dec1_y + diamond_h + gap, arrow_h, arrow_w, "down", GREEN)

    # No arrow up to repair
    add_text(slide, dec1_x + diamond_w + Inches(0.02), dec1_y - Inches(0.02), Inches(0.2), Inches(0.15),
             "No", size=7, bold=True, color=RED)

    # === ROW 3: EXPLAIN and Execute ===
    row3_y = row2_y + box_h + gap + arrow_w + gap + Inches(0.1)

    # 8. EXPLAIN
    explain_x = dec1_x + diamond_w/2 - box_w/2
    add_box(slide, explain_x, row3_y, box_w, box_h, "8. EXPLAIN", "Safe dry-run", TEAL, 9, 7)

    # Arrow left to decision 2
    add_arrow(slide, explain_x - arrow_w - gap, row3_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "left")

    # Decision diamond 2
    dec2_x = explain_x - arrow_w - gap - gap - diamond_w
    dec2_y = row3_y + box_h/2 - diamond_h/2
    add_diamond(slide, dec2_x, dec2_y, diamond_w, diamond_h, "OK?", GRAY)

    # Yes arrow down
    add_text(slide, dec2_x - Inches(0.25), dec2_y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.15),
             "Yes", size=7, bold=True, color=GREEN)
    add_arrow(slide, dec2_x + diamond_w/2 - arrow_h/2, dec2_y + diamond_h + gap, arrow_h, arrow_w, "down", GREEN)

    # No arrow - text indicator
    add_text(slide, dec2_x + diamond_w + Inches(0.02), dec2_y - Inches(0.02), Inches(0.2), Inches(0.15),
             "No", size=7, bold=True, color=RED)

    # === ROW 4: Execute and Return ===
    row4_y = row3_y + box_h + gap + arrow_w + gap + Inches(0.1)

    # 9. Execute
    exec_x = dec2_x + diamond_w/2 - box_w/2
    add_box(slide, exec_x, row4_y, box_w, box_h, "9. Execute", "PostgreSQL", GREEN, 9, 7)

    # Arrow right
    add_arrow(slide, exec_x + box_w + gap, row4_y + box_h/2 - arrow_h/2, arrow_w, arrow_h, "right")

    # 10. Return
    return_x = exec_x + box_w + gap + arrow_w + gap
    add_box(slide, return_x, row4_y, box_w, box_h, "10. Return", "Results", BLUE, 9, 7)

    # === REPAIR LOOP BOX (right side) ===
    repair_x = Inches(9.3)
    repair_y = Inches(0.55)
    repair_w = Inches(2.1)
    repair_h = Inches(1.9)

    repair_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, repair_x, repair_y, repair_w, repair_h)
    repair_shape.fill.solid()
    repair_shape.fill.fore_color.rgb = RGBColor(255, 240, 240)
    repair_shape.line.color.rgb = RED
    repair_shape.line.width = Pt(1.5)

    add_text(slide, repair_x + Inches(0.1), repair_y + Inches(0.08), repair_w - Inches(0.2), Inches(0.2),
             "REPAIR LOOP", size=10, bold=True, color=RED, align="center")
    add_text(slide, repair_x + Inches(0.15), repair_y + Inches(0.35), repair_w - Inches(0.3), Inches(1.4),
             "On error:\n• Collect context\n• Call /repair_sql\n• Regenerate SQL\n• Max 3 attempts\n\n↰ Back to step 4",
             size=8, color=RGBColor(100, 50, 50))

    # Arrow from "No" decisions pointing to repair box
    add_text(slide, dec1_x + diamond_w + Inches(0.15), dec1_y + diamond_h/2 - Inches(0.3), Inches(0.8), Inches(0.2),
             "→ Repair", size=7, italic=True, color=RED)
    add_text(slide, dec2_x + diamond_w + Inches(0.15), dec2_y + diamond_h/2 - Inches(0.3), Inches(0.8), Inches(0.2),
             "→ Repair", size=7, italic=True, color=RED)

    # === SECURITY LAYERS (bottom) ===
    sec_y = Inches(4.05)
    add_text(slide, Inches(0.25), sec_y, Inches(2), Inches(0.2),
             "Security Layers:", size=10, bold=True)

    sec_box_y = sec_y + Inches(0.25)
    sec_box_w = Inches(2.15)
    sec_box_h = Inches(0.5)
    sec_gap = Inches(0.08)

    layers = [
        ("1. Semantic", "Entity match", BLUE),
        ("2. Structural", "SELECT-only", GREEN),
        ("3. EXPLAIN", "Dry-run", TEAL),
        ("4. DB Role", "Permissions", RED),
        ("5. Transaction", "READ ONLY", PURPLE),
    ]

    for i, (title, desc, color) in enumerate(layers):
        lx = Inches(0.25) + i * (sec_box_w + sec_gap)
        add_box(slide, lx, sec_box_y, sec_box_w, sec_box_h, title, desc, color, 8, 7)

    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "NL2SQL MCP Server", "Natural Language to SQL via Model Context Protocol")

    # Slide 2: Problem & Solution
    add_content_slide(prs, "Problem & Solution", [
        "Problem: LLMs generating SQL directly is risky",
        "  SQL injection vulnerabilities",
        "  Hallucinated table/column names",
        "  Dangerous operations (DROP, DELETE)",
        "",
        "Solution: Multi-layer validation pipeline",
        "  Separate SQL generation from execution",
        "  Validate before any database access",
        "  Bounded repair loop for error recovery",
        "  Defense-in-depth security model",
    ])

    # Slide 3: Flowchart
    add_flowchart_slide(prs)

    # Slide 4: Tech Stack
    add_content_slide(prs, "Tech Stack", [
        "Languages:",
        "  TypeScript — MCP server, validation, orchestration",
        "  Python — SQL generation, semantic validation",
        "",
        "Frameworks & Libraries:",
        "  MCP SDK (@modelcontextprotocol/sdk) — Tool registration",
        "  FastAPI — Python REST API server",
        "  pg (node-postgres) — PostgreSQL client",
        "  Zod — Schema validation",
        "",
        "AI/ML:",
        "  Ollama — Local LLM inference server",
        "  HridaAI/hrida-t2sql — Text-to-SQL model (temp=0.0)",
        "",
        "Database: PostgreSQL 15+",
    ])

    # Slide 5: Architecture
    add_content_slide(prs, "Component Details", [
        "LibreChat (Port 3080) — Chat interface",
        "  User types natural language question",
        "",
        "TypeScript MCP Server (stdio) — Orchestration",
        "  Registers nl_query tool, manages validation loop",
        "",
        "Python Sidecar (Port 8001) — SQL Generation",
        "  Filters schema, composes prompt, calls Ollama",
        "",
        "Ollama (Port 11434) — AI Model",
        "  HridaAI/hrida-t2sql, temperature=0.0 (deterministic)",
        "",
        "PostgreSQL (Port 5432) — Database",
        "  MCPtest database, SELECT-only role",
    ])

    # Slide 6: Validation Deep Dive
    add_content_slide(prs, "SQL Validation — How It Works", [
        "State Machine Tokenizer parses SQL safely:",
        "  Distinguishes code from string literals",
        "  Handles 'quotes', \"identifiers\", $$dollar quotes$$",
        "  Ignores -- comments and /* block comments */",
        "",
        "Validation Rules:",
        "  FAIL-FAST: Non-SELECT query → reject immediately",
        "  FAIL-FAST: Multiple statements → reject (SQL injection)",
        "  FAIL-FAST: DROP/DELETE/INSERT → reject",
        "  REPAIRABLE: Unknown table → retry with feedback",
        "  AUTO-FIX: Missing LIMIT → add LIMIT 1000",
    ])

    # Slide 7: Repair Loop Deep Dive
    add_content_slide(prs, "Repair Loop — Error Recovery", [
        "Bounded to 3 attempts (prevents infinite loops)",
        "",
        "Attempt 1: Generate fresh SQL from question",
        "",
        "On validation error (e.g., wrong table name):",
        "  Send error context back to Python sidecar",
        "  Python builds repair prompt with guidance",
        "  Ollama regenerates SQL with corrections",
        "",
        "On PostgreSQL error (e.g., undefined column):",
        "  Parse SQLSTATE error code",
        "  Fail-fast if permission/connection error",
        "  Retry if syntax/column/type error",
        "",
        "Confidence decreases 10% per retry attempt",
    ])

    # Slide 8: All Security Layers
    add_content_slide(prs, "Defense-in-Depth: 5 Security Layers", [
        "Layer 1: Semantic Validation (Python)",
        "  Verifies question entities appear in SQL; detects hallucinations",
        "",
        "Layer 2: Structural Validation (TypeScript)",
        "  Blocks dangerous keywords (DROP, DELETE) and functions (pg_sleep)",
        "",
        "Layer 3: EXPLAIN-First Check",
        "  Safe dry-run catches errors without executing query",
        "",
        "Layer 4: Database Role",
        "  PostgreSQL user has SELECT-only permission",
        "",
        "Layer 5: Transaction Safety",
        "  READ ONLY transactions with statement timeout",
    ], font_size=17)

    # Slide 9: Results
    add_content_slide(prs, "Test Results", [
        "Test Suite: 27 questions across 6 difficulty levels",
        "",
        "Without validation stack: 77.8% success (21/27)",
        "With validation stack: 92.6% success (25/27)",
        "  +14.8% improvement",
        "",
        "By difficulty level:",
        "  Simple queries: 80% → 100%",
        "  Edge cases: 25% → 75%",
        "  NL variations: 80% → 100%",
        "",
        "Latency: ~750ms average (target <3s met)",
        "Security violations: 0",
    ])

    # Slide 10: Summary
    add_content_slide(prs, "Summary", [
        "Architecture eliminates direct LLM-to-database access",
        "",
        "Key benefits:",
        "  Zero SQL injection risk (multi-layer validation)",
        "  Self-healing (bounded repair loop)",
        "  Auditable (every query logged with trace)",
        "",
        "92.6% success rate on test suite",
        "Sub-second latency",
        "Production-ready security model",
    ])

    return prs

def main():
    prs = create_presentation()

    output_path = "/mnt/c/Users/noahc/Downloads/NL2SQL_MCP_Architecture_v7.pptx"
    prs.save(output_path)
    print(f"Saved to: {output_path}")

    project_path = "/home/noahc/nl2sql-project/NL2SQL_MCP_Architecture.pptx"
    prs.save(project_path)
    print(f"Also saved to: {project_path}")

if __name__ == "__main__":
    main()
